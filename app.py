import streamlit as st
import os
import tempfile
import io
import re

# --- 1. 导入基础依赖 ---
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_openai import ChatOpenAI
from langchain_community.vectorstores import Chroma
from langchain_core.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_core.documents import Document  # 用于构建 PPT 文档对象

# --- 2. 导入新功能依赖 (PPT 和 Word) ---
from pptx import Presentation
from docx import Document as DocxDocument
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
# ✨ 关键修复：导入 XML 命名空间处理中文字体
from docx.oxml.ns import qn

# --- 页面配置 ---
st.set_page_config(page_title="DeepSeek 全能研报助手 (修复乱码版)", layout="wide", page_icon="📝")

st.title("📝 DeepSeek 全能研报生成器")
st.markdown("支持 **PDF & PPT** 混合上传 | 生成 **Word (.docx)** 报告 (已修复中文乱码)")
st.markdown("---")

# --- 侧边栏 ---
with st.sidebar:
    st.header("⚙️ 系统设置")
    api_key = st.text_input("请输入 DeepSeek API Key", type="password")
    st.markdown("[👉 点击这里申请 DeepSeek Key](https://platform.deepseek.com/)")
    st.markdown("---")
    st.info("💡 **升级说明**：\n1. 修复了导出 Word 时中文显示为方框的问题。\n2. 默认使用 **微软雅黑** 字体。")


# --- 核心功能函数 ---

@st.cache_resource
def get_embedding_model():
    with st.spinner("正在加载本地向量模型..."):
        return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")


def load_ppt_file(file_path):
    """解析 PPT 文件"""
    prs = Presentation(file_path)
    documents = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            content = "\n".join(slide_text)
            metadata = {"page": i + 1}
            documents.append(Document(page_content=content, metadata=metadata))

    return documents


def process_files(uploaded_files, embeddings):
    """处理混合文件 (PDF + PPT)"""
    if not uploaded_files:
        return None

    all_documents = []
    progress_text = "正在解析文档..."
    my_bar = st.progress(0, text=progress_text)
    total_files = len(uploaded_files)

    for i, uploaded_file in enumerate(uploaded_files):
        my_bar.progress((i / total_files), text=f"正在解析: {uploaded_file.name}")

        file_ext = os.path.splitext(uploaded_file.name)[1].lower()

        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name

        try:
            docs = []
            if file_ext == ".pdf":
                loader = PyPDFLoader(tmp_file_path)
                docs = loader.load()
            elif file_ext in [".ppt", ".pptx"]:
                docs = load_ppt_file(tmp_file_path)
            else:
                st.warning(f"不支持的文件格式: {uploaded_file.name}")
                continue

            for doc in docs:
                doc.metadata['source_filename'] = uploaded_file.name

            all_documents.extend(docs)

        except Exception as e:
            st.error(f"解析 {uploaded_file.name} 失败: {e}")
        finally:
            if os.path.exists(tmp_file_path):
                os.remove(tmp_file_path)

    my_bar.progress(1.0, text="文档解析完成，正在建立向量库...")

    if not all_documents:
        return None

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150
    )
    texts = text_splitter.split_documents(all_documents)

    db = Chroma.from_documents(texts, embeddings)
    my_bar.empty()
    return db


def generate_report(db, topic, api_key):
    """调用 DeepSeek 生成内容"""
    prompt_template = """
    你是一个专业的高级商业分析师。请基于以下【多份文档内容】撰写一份详细的分析报告。

    【综合参考信息】:
    {context}

    【用户指令】: 
    {question}

    【撰写要求】:
    1. **格式**: 必须使用 Markdown 格式（使用 # 表示一级标题，## 表示二级标题，- 表示列表）。
    2. **内容**: 深度整合不同文档的数据。
    3. **引用**: 在关键数据后标注来源。
    4. **严谨**: 仅基于给定材料，不编造。

    请开始撰写:
    """

    PROMPT = PromptTemplate(template=prompt_template, input_variables=["context", "question"])

    llm = ChatOpenAI(
        model_name="deepseek-chat",
        openai_api_key=api_key,
        openai_api_base="https://api.deepseek.com",
        temperature=0,
        max_tokens=3000
    )

    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=db.as_retriever(search_kwargs={"k": 8}),
        chain_type_kwargs={"prompt": PROMPT},
        return_source_documents=True
    )

    return qa_chain.invoke(topic)


def create_word_docx(markdown_text):
    """
    ✨ 修复版：生成支持中文的 Word 文档
    """
    doc = DocxDocument()

    # --- 关键修复：设置全局中文字体 ---
    style = doc.styles['Normal']
    style.font.name = 'Microsoft YaHei'  # 设置西文字体
    style.font.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')  # 设置中文字体 (微软雅黑)
    style.font.size = Pt(11)

    lines = markdown_text.split('\n')

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 简单解析 Markdown 标题
        if line.startswith('# '):
            heading = doc.add_heading(line.replace('# ', ''), level=1)
            # 为标题也设置字体（防止标题乱码）
            for run in heading.runs:
                run.font.name = 'Microsoft YaHei'
                run.font.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

        elif line.startswith('## '):
            heading = doc.add_heading(line.replace('## ', ''), level=2)
            for run in heading.runs:
                run.font.name = 'Microsoft YaHei'
                run.font.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

        elif line.startswith('### '):
            heading = doc.add_heading(line.replace('### ', ''), level=3)
            for run in heading.runs:
                run.font.name = 'Microsoft YaHei'
                run.font.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

        elif line.startswith('- ') or line.startswith('* '):
            p = doc.add_paragraph(line.replace('- ', '').replace('* ', ''), style='List Bullet')
        else:
            doc.add_paragraph(line)

    # 保存到内存流
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio


# --- 主界面逻辑 ---

if not api_key:
    st.warning("⚠️ 请先在左侧侧边栏输入 DeepSeek API Key。")
else:
    embedding_model = get_embedding_model()

    uploaded_files = st.file_uploader(
        "📄 上传文档 (支持 PDF 和 PPTX)",
        type=["pdf", "pptx", "ppt"],
        accept_multiple_files=True
    )

    if uploaded_files:
        current_file_names = [f.name for f in uploaded_files]

        if "last_uploaded_files_mix" not in st.session_state or st.session_state.last_uploaded_files_mix != current_file_names:
            st.session_state.vector_db_mix = process_files(uploaded_files, embedding_model)
            st.session_state.last_uploaded_files_mix = current_file_names
            if st.session_state.vector_db_mix:
                st.success(f"✅ 已成功解析 {len(uploaded_files)} 份文档！")

        st.subheader("📊 报告生成设置")
        default_topic = "综合分析这些文档，输出一份包含摘要、关键发现和结论的完整报告。"
        report_topic = st.text_area("分析指令:", value=default_topic, height=100)

        if st.button("🚀 生成并导出报告"):
            if "vector_db_mix" in st.session_state:
                with st.spinner("正在思考并撰写 Word 报告..."):
                    try:
                        response = generate_report(st.session_state.vector_db_mix, report_topic, api_key)
                        report_content = response['result']

                        st.markdown("### 📄 报告预览")
                        st.markdown(report_content)
                        st.markdown("---")

                        # 生成修复乱码后的 Word
                        docx_file = create_word_docx(report_content)

                        st.download_button(
                            label="📥 下载 Word 报告 (.docx)",
                            data=docx_file,
                            file_name="DeepSeek_分析报告.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )

                        with st.expander("🔎 查看引用来源"):
                            for i, doc in enumerate(response['source_documents']):
                                source = doc.metadata.get('source_filename', '未知文件')
                                page = doc.metadata.get('page', '?')
                                st.markdown(f"**[{i + 1}] {source} (第 {page} 页/张):**")
                                st.caption(f"> {doc.page_content[:150]}...")
                                st.divider()

                    except Exception as e:
                        st.error(f"生成失败: {e}")
            else:
                st.error("请等待文档解析完成。")